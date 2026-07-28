# -*- coding: utf-8 -*-
"""Builds Documind_AI_Presentation.pptx — enterprise dark theme matching the
product's own emerald-on-black branding.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import deck_content as C

EMU_PER_IN = 914400
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Segoe UI"


def rgb(hex_str):
    return RGBColor.from_string(hex_str)


def set_bg(slide, color=C.BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(color)


def no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    if line:
        shp.line.color.rgb = rgb(C.BORDER)
        shp.line.width = Pt(0.75)
    else:
        no_line(shp)
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, color, radius=0.08, line=False, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    if line:
        shp.line.color.rgb = rgb(line_color or C.BORDER)
        shp.line.width = Pt(1)
    else:
        no_line(shp)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=C.TEXT, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = rgb(color)
    return box


def add_kicker(slide, x, y, text, color=C.ACCENT):
    return add_text(slide, x, y, Inches(8), Inches(0.35), text.upper(), size=13, color=color, bold=True, font=FONT)


def add_footer(slide, page_no, total, section=""):
    add_rect(slide, 0, SLIDE_H - Inches(0.04), SLIDE_W, Inches(0.04), C.ACCENT)
    add_text(slide, Inches(0.55), SLIDE_H - Inches(0.45), Inches(6), Inches(0.3),
              "DOCUMIND AI", size=9, color=C.MUTED, bold=True)
    if section:
        add_text(slide, Inches(4.5), SLIDE_H - Inches(0.45), Inches(5), Inches(0.3),
                  section.upper(), size=9, color=C.MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45), Inches(0.8), Inches(0.3),
              f"{page_no:02d} / {total:02d}", size=9, color=C.MUTED, align=PP_ALIGN.RIGHT)


def _set_alpha(shape, pct_visible):
    """pct_visible: 0-100, how opaque the fill is (100 = fully solid)."""
    srgb = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(pct_visible * 1000))})
    srgb.append(alpha)


def glow(slide, cx, cy, r, color=C.ACCENT):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    _set_alpha(shp, 12)
    no_line(shp)
    shp.shadow.inherit = False
    return shp


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    return slide


def section_header(slide, kicker, title, subtitle=None):
    add_kicker(slide, Inches(0.7), Inches(0.55), kicker)
    add_text(slide, Inches(0.7), Inches(0.85), Inches(11.5), Inches(0.9), title,
              size=32, bold=True, color=C.TEXT)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(1.55), Inches(11.5), Inches(0.5), subtitle,
                  size=14, color=C.MUTED)


# ---------------------------------------------------------------- slides ---

def slide_title(prs):
    slide = new_slide(prs)
    glow(slide, Inches(-2), Inches(-2), Inches(9), C.ACCENT)
    glow(slide, Inches(9), Inches(4), Inches(7), C.BLUE)
    add_rect(slide, Inches(0.7), Inches(2.55), Inches(0.55), Inches(0.08), C.ACCENT)
    add_text(slide, Inches(0.7), Inches(2.75), Inches(11.5), Inches(1.3), C.TITLE,
              size=56, bold=True, color=C.TEXT)
    add_text(slide, Inches(0.72), Inches(3.7), Inches(10.5), Inches(0.7), C.SUBTITLE,
              size=18, color=C.MUTED)
    add_text(slide, Inches(0.72), Inches(5.6), Inches(8), Inches(0.4), C.AUTHOR, size=16, bold=True, color=C.TEXT)
    add_text(slide, Inches(0.72), Inches(6.0), Inches(8), Inches(0.35), C.UNIVERSITY, size=13, color=C.MUTED)
    add_text(slide, Inches(0.72), Inches(6.3), Inches(8), Inches(0.35), C.DEPARTMENT, size=13, color=C.MUTED)
    add_text(slide, Inches(0.72), Inches(6.6), Inches(8), Inches(0.35), C.SUPERVISOR, size=13, color=C.MUTED)
    add_text(slide, SLIDE_W - Inches(2.2), Inches(6.6), Inches(1.6), Inches(0.35), C.DATE,
              size=13, color=C.ACCENT, bold=True, align=PP_ALIGN.RIGHT)


def slide_agenda(prs, idx, total):
    slide = new_slide(prs)
    section_header(slide, "Overview", "Agenda")
    y = Inches(2.1)
    for i, item in enumerate(C.AGENDA):
        row_y = y + Inches(0.55) * i
        add_rounded(slide, Inches(0.7), row_y, Inches(0.5), Inches(0.4), C.SURFACE, radius=0.5, line=True)
        add_text(slide, Inches(0.7), row_y + Inches(0.03), Inches(0.5), Inches(0.4), f"{i+1:02d}",
                  size=13, bold=True, color=C.ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(1.4), row_y + Inches(0.03), Inches(9), Inches(0.4), item,
                  size=16, color=C.TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, idx, total)


def slide_problem(prs, idx, total):
    slide = new_slide(prs)
    section_header(slide, "The Problem", "Working with PDFs is Still Slow and Fragmented")
    cols = 2
    card_w, card_h = Inches(5.55), Inches(2.0)
    gap_x, gap_y = Inches(0.4), Inches(0.35)
    start_x, start_y = Inches(0.7), Inches(2.25)
    for i, (title, body) in enumerate(C.PROBLEM_POINTS):
        col, row = i % cols, i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_rounded(slide, x, y, card_w, card_h, C.SURFACE, radius=0.06, line=True)
        add_rect(slide, x, y, Inches(0.06), card_h, C.ACCENT)
        add_text(slide, x + Inches(0.35), y + Inches(0.25), card_w - Inches(0.7), Inches(0.4),
                  title, size=17, bold=True, color=C.TEXT)
        add_text(slide, x + Inches(0.35), y + Inches(0.75), card_w - Inches(0.7), Inches(1.1),
                  body, size=13, color=C.MUTED, line_spacing=1.15)
    add_footer(slide, idx, total, "Problem Statement")


def slide_solution(prs, idx, total):
    slide = new_slide(prs)
    section_header(slide, "The Solution", "Documind AI: One Platform, Every Document Workflow")
    y = Inches(2.3)
    for i, point in enumerate(C.SOLUTION_POINTS):
        row_y = y + Inches(0.9) * i
        add_rounded(slide, Inches(0.7), row_y, Inches(0.5), Inches(0.5), f"{C.ACCENT}", radius=0.5)
        add_text(slide, Inches(0.7), row_y, Inches(0.5), Inches(0.5), "✓", size=18, bold=True,
                  color=C.BG, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(1.45), row_y - Inches(0.03), Inches(10.8), Inches(0.8), point,
                  size=16, color=C.TEXT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    add_footer(slide, idx, total, "Solution Overview")


def slide_tech_stack(prs, idx, total):
    slide = new_slide(prs)
    section_header(slide, "Technology", "Technology Stack")
    cols = 2
    card_w, card_h = Inches(5.55), Inches(1.05)
    gap_x, gap_y = Inches(0.4), Inches(0.22)
    start_x, start_y = Inches(0.7), Inches(2.15)
    for i, (label, val) in enumerate(C.TECH_STACK):
        col, row = i % cols, i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_rounded(slide, x, y, card_w, card_h, C.SURFACE, radius=0.12, line=True)
        add_text(slide, x + Inches(0.3), y + Inches(0.12), card_w - Inches(0.6), Inches(0.3),
                  label.upper(), size=11, bold=True, color=C.ACCENT)
        add_text(slide, x + Inches(0.3), y + Inches(0.45), card_w - Inches(0.6), Inches(0.5),
                  val, size=12.5, color=C.TEXT)
    add_footer(slide, idx, total, "Technology Stack")


def slide_architecture(prs, idx, total):
    slide = new_slide(prs)
    section_header(slide, "System Design", "System Architecture")
    x = Inches(0.9)
    w = Inches(11.5)
    y = Inches(2.2)
    h = Inches(0.85)
    gap = Inches(0.15)
    for label, desc, color in C.ARCHITECTURE_LAYERS:
        add_rounded(slide, x, y, w, h, C.SURFACE, radius=0.1, line=True, line_color=color)
        add_rect(slide, x, y, Inches(0.08), h, color)
        add_text(slide, x + Inches(0.35), y + Inches(0.12), Inches(2.3), Inches(0.6), label,
                  size=15, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(2.8), y + Inches(0.12), w - Inches(3.1), Inches(0.6), desc,
                  size=12.5, color=C.MUTED, anchor=MSO_ANCHOR.MIDDLE)
        y += h + gap
    add_footer(slide, idx, total, "System Architecture")


def slide_core_features(prs, idx, total):
    slide = new_slide(prs)
    section_header(slide, "Product", "Core Features")
    cols = 3
    card_w, card_h = Inches(3.65), Inches(1.55)
    gap_x, gap_y = Inches(0.3), Inches(0.25)
    start_x, start_y = Inches(0.7), Inches(2.05)
    for i, (title, body) in enumerate(C.CORE_FEATURES):
        col, row = i % cols, i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_rounded(slide, x, y, card_w, card_h, C.SURFACE, radius=0.07, line=True)
        add_text(slide, x + Inches(0.25), y + Inches(0.18), card_w - Inches(0.5), Inches(0.35),
                  title, size=14, bold=True, color=C.TEXT)
        add_text(slide, x + Inches(0.25), y + Inches(0.55), card_w - Inches(0.5), Inches(0.9),
                  body, size=10.5, color=C.MUTED, line_spacing=1.1)
    add_footer(slide, idx, total, "Core Features")


def slide_standout(prs, idx, total):
    slide = new_slide(prs)
    section_header(slide, "What Makes It Different", "Standout Innovations")
    colors = [C.ACCENT, C.BLUE, C.PURPLE, C.ACCENT, C.AMBER]
    y = Inches(2.1)
    for i, (title, body) in enumerate(C.STANDOUT_FEATURES):
        row_h = Inches(0.92)
        color = colors[i % len(colors)]
        add_rounded(slide, Inches(0.7), y, Inches(11.9), row_h, C.SURFACE, radius=0.08, line=True, line_color=color)
        add_rect(slide, Inches(0.7), y, Inches(0.08), row_h, color)
        add_text(slide, Inches(1.05), y + Inches(0.1), Inches(3.1), Inches(0.7), title,
                  size=14.5, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(4.3), y + Inches(0.1), Inches(8.1), Inches(0.75), body,
                  size=11.5, color=C.MUTED, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        y += row_h + Inches(0.12)
    add_footer(slide, idx, total, "Standout Innovations")


def slide_security(prs, idx, total):
    slide = new_slide(prs)
    section_header(slide, "Trust", "Security, Privacy & Billing")
    y = Inches(2.15)
    for title, body in C.SECURITY_POINTS:
        row_h = Inches(0.85)
        add_rounded(slide, Inches(0.7), y, Inches(11.9), row_h, C.SURFACE, radius=0.08, line=True)
        add_text(slide, Inches(1.0), y + Inches(0.1), Inches(3.0), Inches(0.65), title,
                  size=13.5, bold=True, color=C.ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(4.1), y + Inches(0.1), Inches(8.3), Inches(0.65), body,
                  size=11, color=C.MUTED, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        y += row_h + Inches(0.12)
    add_footer(slide, idx, total, "Security & Privacy")


def slide_challenges(prs, idx, total):
    slide = new_slide(prs)
    section_header(slide, "Engineering", "Challenges & Key Learnings")
    cols = 2
    card_w, card_h = Inches(5.55), Inches(2.15)
    gap_x, gap_y = Inches(0.4), Inches(0.3)
    start_x, start_y = Inches(0.7), Inches(2.2)
    for i, (title, body) in enumerate(C.CHALLENGES):
        col, row = i % cols, i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_rounded(slide, x, y, card_w, card_h, C.SURFACE, radius=0.06, line=True)
        add_rect(slide, x, y, Inches(0.06), card_h, C.AMBER)
        add_text(slide, x + Inches(0.35), y + Inches(0.22), card_w - Inches(0.7), Inches(0.5),
                  title, size=15, bold=True, color=C.TEXT)
        add_text(slide, x + Inches(0.35), y + Inches(0.75), card_w - Inches(0.7), Inches(1.3),
                  body, size=11.5, color=C.MUTED, line_spacing=1.15)
    add_footer(slide, idx, total, "Challenges & Learnings")


def slide_demo(prs, idx, total):
    slide = new_slide(prs)
    glow(slide, Inches(4), Inches(2), Inches(8), C.ACCENT)
    add_kicker(slide, Inches(0.7), Inches(2.5), "Live Demonstration")
    add_text(slide, Inches(0.7), Inches(2.85), Inches(11.5), Inches(1.0), "Let's See It in Action",
              size=40, bold=True, color=C.TEXT)
    add_text(slide, Inches(0.72), Inches(3.85), Inches(10), Inches(0.5),
              "Chat → Audio Overview → Knowledge Graph → Live Collaboration",
              size=16, color=C.MUTED)
    add_rounded(slide, Inches(0.72), Inches(4.6), Inches(5.4), Inches(0.65), C.SURFACE, radius=0.5, line=True)
    add_text(slide, Inches(0.72), Inches(4.6), Inches(5.4), Inches(0.65), C.LIVE_URL,
              size=15, bold=True, color=C.ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, idx, total, "Live Demo")


def slide_thanks(prs, idx, total):
    slide = new_slide(prs)
    glow(slide, Inches(-1), Inches(4), Inches(7), C.BLUE)
    glow(slide, Inches(9), Inches(-1), Inches(7), C.ACCENT)
    add_text(slide, Inches(0.7), Inches(2.6), Inches(11), Inches(1.2), "Thank You",
              size=52, bold=True, color=C.TEXT)
    add_text(slide, Inches(0.72), Inches(3.7), Inches(10), Inches(0.5), "Questions & Discussion",
              size=18, color=C.MUTED)
    y = Inches(4.7)
    for label, url in C.CLOSING_LINKS:
        add_text(slide, Inches(0.72), y, Inches(2.2), Inches(0.35), label.upper(),
                  size=11, bold=True, color=C.ACCENT)
        add_text(slide, Inches(0.72), y + Inches(0.32), Inches(9), Inches(0.4), url,
                  size=14, color=C.TEXT)
        y += Inches(0.75)
    add_footer(slide, idx, total)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        ("title", slide_title),
        ("agenda", slide_agenda),
        ("problem", slide_problem),
        ("solution", slide_solution),
        ("stack", slide_tech_stack),
        ("architecture", slide_architecture),
        ("core", slide_core_features),
        ("standout", slide_standout),
        ("security", slide_security),
        ("challenges", slide_challenges),
        ("demo", slide_demo),
        ("thanks", slide_thanks),
    ]
    total = len(builders)

    slide_title(prs)
    idx = 2
    for name, fn in builders[1:]:
        fn(prs, idx, total)
        idx += 1

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "Documind_AI_Presentation.pptx")
    prs.save(out_path)
    print("Saved:", out_path)


if __name__ == "__main__":
    main()
